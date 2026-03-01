#!/usr/bin/env python3
"""
generate_pptx.py
-----------------
Generates presentation.pptx for the Brain Tumor Segmentation FPGA project.
Content mirrors presentation_outline.md.

Usage:
    python generate_pptx.py

Output:
    presentation.pptx  (in the same directory)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import os

OUT_DIR = os.path.dirname(os.path.abspath(__file__))
OUT_PATH = os.path.join(OUT_DIR, "presentation.pptx")

# Colors
TITLE_COLOR = RGBColor(0x1B, 0x3A, 0x5C)   # dark navy
ACCENT_COLOR = RGBColor(0x2E, 0x86, 0xC1)  # blue accent
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
LIGHT_BG = RGBColor(0xEC, 0xF0, 0xF1)
GREEN = RGBColor(0x27, 0xAE, 0x60)


def set_slide_bg(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, TITLE_COLOR)

    # Main title
    txBox = slide.shapes.add_textbox(
        Inches(0.8), Inches(1.5), Inches(8.4), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Brain Tumor Segmentation on FPGA"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    txBox2 = slide.shapes.add_textbox(
        Inches(0.8), Inches(3.2), Inches(8.4), Inches(0.8))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "with Adaptive Processing Mode Selection"
    p2.font.size = Pt(22)
    p2.font.italic = True
    p2.font.color.rgb = ACCENT_COLOR
    p2.alignment = PP_ALIGN.CENTER

    # Details
    details = [
        "Digital System Design Lab",
        "Board: Nexys A7-100T (Artix-7 xc7a100tcsg324-1)",
        "Tools: Vitis HLS / Vivado 2025.1",
    ]
    txBox3 = slide.shapes.add_textbox(
        Inches(1.5), Inches(4.5), Inches(7), Inches(2))
    tf3 = txBox3.text_frame
    tf3.word_wrap = True
    for d in details:
        p3 = tf3.add_paragraph(
        ) if tf3.paragraphs[0].text else tf3.paragraphs[0]
        p3.text = d
        p3.font.size = Pt(16)
        p3.font.color.rgb = RGBColor(0xBD, 0xC3, 0xC7)
        p3.alignment = PP_ALIGN.CENTER
        p3 = tf3.add_paragraph()


def add_content_slide(prs, title, bullets, sub_bullets=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, WHITE)

    # Title bar
    shape = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(10), Inches(1.1))  # rectangle
    shape.fill.solid()
    shape.fill.fore_color.rgb = TITLE_COLOR
    shape.line.fill.background()

    txBox = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.15), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Body bullets
    txBox2 = slide.shapes.add_textbox(
        Inches(0.6), Inches(1.4), Inches(8.8), Inches(5.5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True

    for i, b in enumerate(bullets):
        p2 = tf2.add_paragraph(
        ) if i > 0 or tf2.paragraphs[0].text else tf2.paragraphs[0]
        p2.text = b
        p2.font.size = Pt(18)
        p2.font.color.rgb = DARK_GRAY
        p2.space_after = Pt(8)
        p2.level = 0

        if sub_bullets and i in sub_bullets:
            for sb in sub_bullets[i]:
                p3 = tf2.add_paragraph()
                p3.text = sb
                p3.font.size = Pt(15)
                p3.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
                p3.space_after = Pt(4)
                p3.level = 1


def add_table_slide(prs, title, headers, rows):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    # Title bar
    shape = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(10), Inches(1.1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = TITLE_COLOR
    shape.line.fill.background()

    txBox = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.15), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Table
    num_rows = len(rows) + 1
    num_cols = len(headers)
    left = Inches(0.8)
    top = Inches(1.6)
    width = Inches(8.4)
    height = Inches(0.5 * num_rows)

    table_shape = slide.shapes.add_table(
        num_rows, num_cols, left, top, width, height)
    table = table_shape.table

    # Header row
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT_COLOR
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(14)
            paragraph.font.bold = True
            paragraph.font.color.rgb = WHITE
            paragraph.alignment = PP_ALIGN.CENTER

    # Data rows
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if i % 2 == 0 else LIGHT_BG
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(13)
                paragraph.font.color.rgb = DARK_GRAY
                paragraph.alignment = PP_ALIGN.CENTER


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    add_title_slide(prs)

    # Slide 2: Problem Statement
    add_content_slide(prs, "Problem Statement", [
        "Brain tumor segmentation from MRI is critical for diagnosis",
        "Software-only approaches are slow and power-hungry",
        "Goal: FPGA acceleration for real-time, energy-efficient processing",
        "Input: 256×256 grayscale MRI images",
        "Target: Nexys A7-100T (Artix-7, 100 MHz)",
    ])

    # Slide 3: Novelty
    add_table_slide(prs, "Novelty – Adaptive Mode Selection",
                    ["Image Quality", "Mode", "Behaviour"],
                    [
                        ["High contrast", "FAST", "Minimal processing, max speed"],
                        ["Medium contrast", "NORMAL", "Standard Otsu + cleanup"],
                        ["Low contrast", "CAREFUL",
                            "Adaptive threshold + full morph"],
                    ])

    # Slide 4: System Architecture
    add_content_slide(prs, "System Architecture", [
        "MicroBlaze CPU ←AXI→ HLS Otsu Accelerator",
        "Shared 128 KB Image BRAM (dual-port, dual AXI BRAM controllers)",
        "AXI SmartConnect merges HLS gmem0/gmem1 → BRAM",
        "Peripherals: UART (115200), GPIO (5 LEDs), Timer",
        "MicroBlaze controls pipeline; HLS does heavy computation",
    ])

    # Slide 5: Algorithm
    add_content_slide(prs, "Algorithm Pipeline", [
        "1. Histogram – 256-bin grayscale histogram (II=1, pipelined)",
        "2. Otsu Threshold – maximise inter-class variance",
        "3. Morphology – erode/dilate for noise removal (mode-dependent)",
        "4. Adaptive Mode – auto-select FAST/NORMAL/CAREFUL from image stats",
        "5. Watershed – connected-component labelling (software post-processing)",
    ])

    # Slide 6: Implementation Flow
    add_content_slide(prs, "Implementation Flow", [
        "Phase 1: Python – Verify algorithm (Otsu + Watershed + metrics)",
        "Phase 2: HLS C++ – Synthesise Otsu to IP core (xc7a100tcsg324-1)",
        "Phase 3: Vivado – Build MicroBlaze SoC + bitstream",
        "Phase 4: Vitis – MicroBlaze bare-metal firmware",
        "Phase 5: Test Images – Convert PNG → .bin / .h",
        "Phase 6: Documentation – Report, presentation, user manual",
    ])

    # Slide 7: HLS Results
    add_table_slide(prs, "HLS Synthesis Results",
                    ["Step", "Status", "Details"],
                    [
                        ["C Simulation", "✅ PASSED", "All 9 tests, all 3 modes"],
                        ["C Synthesis", "✅ 97 MHz", "II=1 on all major loops"],
                        ["IP Export", "✅ Generated",
                            "custom_hls_otsu_threshold_top_1_0.zip"],
                    ])

    # Slide 8: Resource Usage
    add_table_slide(prs, "FPGA Resource Usage (Estimated)",
                    ["Resource", "Used", "Available", "Utilisation"],
                    [
                        ["LUT", "~7,600", "63,400", "~12%"],
                        ["FF", "~3,800", "126,800", "~3%"],
                        ["BRAM", "5.5", "135", "~4%"],
                        ["DSP", "7", "240", "~3%"],
                    ])

    # Slide 9: Segmentation Results
    add_table_slide(prs, "Segmentation Results",
                    ["Image", "Dice Score", "Mode Selected"],
                    [
                        ["brain_01 (bright tumor)", "0.98", "FAST"],
                        ["brain_02 (subtle tumor)", "0.98", "NORMAL"],
                        ["brain_03 (no tumor)", "0.19", "CAREFUL"],
                    ])

    # Slide 10: Performance Comparison
    add_table_slide(prs, "Performance Comparison: SW vs FPGA",
                    ["Metric", "SW (Python/CPU)", "HW (FPGA)", "Improvement"],
                    [
                        ["Processing time", "~18 ms", "~2 ms", "8.9× speedup"],
                        ["Power", "15 W", "0.7 W", "21× lower"],
                        ["Energy per image", "~270 mJ", "~1.4 mJ", ">99% savings"],
                    ])

    # Slide 11: Demo
    add_content_slide(prs, "Demo", [
        "Live demo on Nexys A7-100T board",
        "Serial terminal (115200 baud) shows processing results",
        "LEDs indicate current processing mode and status:",
    ], sub_bullets={
        2: [
            "LD0: Heartbeat",
            "LD1: Processing active",
            "LD2-LD3: Mode (00=FAST, 01=NORMAL, 10=CAREFUL)",
            "LD4: Done",
        ]
    })

    # Slide 12: Conclusion
    add_content_slide(prs, "Conclusion", [
        "✅ Successfully implemented Otsu segmentation on FPGA",
        "✅ Adaptive mode selection – novel contribution",
        "✅ ~8.9× speedup, >99% energy savings vs software",
        "✅ Complete pipeline: Python → HLS → Vivado → Vitis",
        "",
        "Future Work:",
    ], sub_bullets={
        5: [
            "Higher resolution (512×512+)",
            "DMA-based data transfer for faster I/O",
            "Real MRI DICOM image support",
            "Zynq port with ARM + Linux",
        ]
    })

    # Slide 13: Q&A
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, TITLE_COLOR)

    txBox = slide.shapes.add_textbox(
        Inches(1), Inches(2.5), Inches(8), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Thank You!"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    txBox2 = slide.shapes.add_textbox(
        Inches(1), Inches(4.2), Inches(8), Inches(1))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "Questions?"
    p2.font.size = Pt(28)
    p2.font.color.rgb = ACCENT_COLOR
    p2.alignment = PP_ALIGN.CENTER

    txBox3 = slide.shapes.add_textbox(
        Inches(1), Inches(5.5), Inches(8), Inches(0.5))
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "github.com/akulasahasra75/FPGA--Brain-Tumor-Segmentation"
    p3.font.size = Pt(14)
    p3.font.color.rgb = RGBColor(0xBD, 0xC3, 0xC7)
    p3.alignment = PP_ALIGN.CENTER

    # Save
    prs.save(OUT_PATH)
    print(f"Presentation saved: {OUT_PATH}")
    print(f"  {len(prs.slides)} slides generated")


if __name__ == "__main__":
    main()
